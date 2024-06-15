from django.shortcuts import render, redirect, get_object_or_404
from django.http import HttpResponse, JsonResponse
from django import forms
from .models import Inventory, Weather, Movie
import csv
import os
import sqlite3
from django.conf import settings
from .forms import QueryForm, UploadCSVForm
from transformers import AutoModelForSeq2SeqLM, AutoTokenizer
from pptx import Presentation
from pptx.util import Inches
from docx import Document
import json
import xml.etree.ElementTree as ET
from .models import QueryHistory

# Define InventoryForm class
class InventoryForm(forms.ModelForm):
    class Meta:
        model = Inventory
        fields = ['product_name', 'quantity_in_stock', 'cost_per_item', 'quantity_sold', 'sales', 'stock_date', 'photos']
        widgets = {
            'stock_date': forms.DateInput(attrs={'type': 'date'}),
        }

# Define QueryForm class
class QueryForm(forms.Form):
    query = forms.CharField(label='Enter your natural language query', max_length=200)

# Load the ML model and tokenizer
model_path = 'gaussalgo/T5-LM-Large-text2sql-spider'
model = AutoModelForSeq2SeqLM.from_pretrained(model_path)
tokenizer = AutoTokenizer.from_pretrained(model_path)


import pandas as pd

def upload_csv(request):
    if request.method == 'POST':
        form = UploadCSVForm(request.POST, request.FILES)
        if form.is_valid():
            csv_file = request.FILES['csv_file']
            file_path = os.path.join(settings.MEDIA_ROOT, csv_file.name)
            
            # Save the uploaded CSV file to a specific location
            with open(file_path, 'wb+') as destination:
                for chunk in csv_file.chunks():
                    destination.write(chunk)
            
            # Read the CSV file using pandas
            try:
                df = pd.read_csv(file_path)
                table_name = os.path.splitext(csv_file.name)[0]

                # Connect to SQLite database
                db_path = os.path.join(settings.BASE_DIR, 'db.sqlite3')
                conn = sqlite3.connect(db_path)

                # Write DataFrame to SQLite database
                df.to_sql(table_name, conn, if_exists='replace', index=False)

                # Close database connection
                conn.close()

                # Set success flag in session
                request.session['upload_success'] = True  # This line ensures upload success flag is set

                # Optionally, you can also refresh the table names in the session or context
                table_names = get_table_names()  # Replace with actual function to fetch table names
                request.session['table_names'] = table_names  # Update table names in session

            except Exception as e:
                print(f"Error processing CSV file: {e}")
                # Handle any exceptions or errors that occur during processing
                request.session['upload_success'] = False  # Ensure to set to False if upload fails

            return redirect('query_view')  # Redirect to your query view after successful upload
    else:
        form = UploadCSVForm()
    
    return render(request, 'upload_csv.html', {'form': form})


def get_table_names():
    db_path = os.path.join(settings.BASE_DIR, 'db.sqlite3')
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()
    cursor.execute("SELECT name FROM sqlite_master WHERE type='table' AND name NOT LIKE 'auth_%' AND name NOT LIKE 'django_%' AND name NOT LIKE 'sqlite%'")
    tables = cursor.fetchall()
    table_names = [table[0] for table in tables]
    conn.close()
    return table_names

def query_database(db_path, query):
    try:
        conn = sqlite3.connect(db_path)
        cursor = conn.cursor()
        cursor.execute(query)
        results = cursor.fetchall()
        column_names = [description[0] for description in cursor.description]
        conn.close()
        return [dict(zip(column_names, row)) for row in results]
    except sqlite3.Error as e:
        return str(e)

def get_sql_query(question, schema):
    input_text = f"Question: {question} Schema: {schema}"
    model_inputs = tokenizer(input_text, return_tensors="pt")
    outputs = model.generate(**model_inputs, max_length=512)
    output_text = tokenizer.batch_decode(outputs, skip_special_tokens=True)
    return output_text[0]

def download_csv(results):
    response = HttpResponse(content_type='text/csv')
    response['Content-Disposition'] = 'attachment; filename="results.csv"'
    writer = csv.writer(response)
    if results:
        headers = results[0].keys()
        writer.writerow(headers)
        for row in results:
            writer.writerow(row.values())
    return response

def download_word(results):
    document = Document()
    document.add_heading('Query Results', 0)
    if results:
        headers = results[0].keys()
        table = document.add_table(rows=1, cols=len(headers))
        hdr_cells = table.rows[0].cells
        for i, header in enumerate(headers):
            hdr_cells[i].text = header
        for row in results:
            row_cells = table.add_row().cells
            for i, value in enumerate(row.values()):
                row_cells[i].text = str(value)
    document_path = 'query_results.docx'
    document.save(document_path)
    with open(document_path, 'rb') as docx_file:
        response = HttpResponse(docx_file.read(), content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
        response['Content-Disposition'] = 'attachment; filename=query_results.docx'
    return response

def download_ppt(results):
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Query Results"
    if not results:
        return HttpResponse("No results found")
    headers = results[0].keys()
    table = slide.shapes.add_table(rows=len(results) + 1, cols=len(headers), left=Inches(0.5), top=Inches(1.5), width=Inches(9), height=Inches(0.8)).table
    for i, header in enumerate(headers):
        table.cell(0, i).text = header
    for i, row in enumerate(results):
        for j, value in enumerate(row.values()):
            table.cell(i + 1, j).text = str(value)
    response = HttpResponse(content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    response['Content-Disposition'] = 'attachment; filename="results.pptx"'
    prs.save(response)
    return response

def download_history_xml(request):
    # query_history = request.session.get('query_history', [])
    user_id = request.user.username
    query_history = QueryHistory.objects.filter(user_id=user_id).values('query_history')
    if query_history:
        query_history = query_history[0]['query_history']
    else:
        query_history = []
    root = ET.Element('QueryHistory')
    for item in query_history:
        entry = ET.SubElement(root, 'Entry')
        question = ET.SubElement(entry, 'Question')
        question.text = item['question']
        sql_query = ET.SubElement(entry, 'SQLQuery')
        sql_query.text = item['sql_query']
        results = ET.SubElement(entry, 'Results')
        if isinstance(item['results'], str):
            result_entry = ET.SubElement(results, 'Result')
            result_field = ET.SubElement(result_entry, 'Error')
            result_field.text = item['results']
        else:
            for result in item['results']:
                result_entry = ET.SubElement(results, 'Result')
                for key, value in result.items():
                    result_field = ET.SubElement(result_entry, key)
                    result_field.text = str(value)
    xml_str = ET.tostring(root, encoding='utf-8')
    response = HttpResponse(xml_str, content_type='application/xml')
    response['Content-Disposition'] = 'attachment; filename=query_history.xml'
    return response

def clear_history(request):
    # if 'query_history' in request.session:
        # del request.session['query_history']
    # delete all query history where user_id = user_id
    user_id = request.user.username
    QueryHistory.objects.filter(user_id=user_id).delete()
    return redirect('query_view')

def query_view(request):
    name = request.GET.get('name', 'default_name')
    results = None
    sql_query = None
    query_history = []
    user_id = request.user.username
    query_history = QueryHistory.objects.filter(user_id=user_id).values('query_history')
    if query_history:
        query_history = query_history[0]['query_history']
    else:
        query_history = []
    print(query_history)


    # Fetch table names, excluding Django inbuilt tables
    db_path = os.path.join(settings.BASE_DIR, 'db.sqlite3')
    conn = sqlite3.connect(db_path)
    cursor = conn.cursor()
    cursor.execute("SELECT name FROM sqlite_master WHERE type='table' AND name NOT LIKE 'auth_%' AND name NOT LIKE 'django_%' AND name NOT LIKE 'sqlite%'")
    tables = cursor.fetchall()
    table_names = [table[0] for table in tables]
    conn.close()

    # Schema definition for custom tables
    schema = {
        'myapps_inventory': '''"myapps_inventory" "DATABASE" 
                            "product_ID" int, 
                            "product_name" text, 
                            "quantity_in_stock" int, 
                            "cost_per_item" float, 
                            primary key: "product_ID"''',
        'myapps_weather': '''"myapps_weather" "DATABASE" 
                            "id" int, 
                            "formatted_date" datetime, 
                            "summary" varchar(200), 
                            "precip_type" varchar(50), 
                            "temperature_c" float, 
                            "apparent_temperature_c" float, 
                            "humidity" float, 
                            "wind_speed_kmh" float, 
                            "wind_bearing_degrees" int, 
                            "visibility_km" float, 
                            "cloud_cover" float, 
                            "pressure_millibars" float, 
                            "daily_summary" text, 
                            primary key: "id"''',
        'myapps_movie': '''"myapps_movie" "DATABASE" 
                            "id" int, 
                            "poster_link" varchar(200), 
                            "series_title" varchar(100), 
                            "released_year" varchar(10), 
                            "certificate" varchar(10), 
                            "runtime" varchar(20), 
                            "genre" varchar(100), 
                            "imdb_rating" float, 
                            "overview" text,
                            "meta_score" varchar(10), 
                            "director" varchar(100), 
                            "star1" varchar(100), 
                            "star2" varchar(100), 
                            "star3" varchar(100), 
                            "star4" varchar(100), 
                            "no_of_votes" bigint, 
                            "gross" varchar(50), 
                            primary key: "id"''',
        'myapps_electronic': '''"myapps_electronic" "DATABASE" 
                            "series_reference" varchar(100), 
                            "period" varchar(10), 
                            "data_value" float, 
                            "suppressed" varchar(1), 
                            "status" varchar(1), 
                            "units" varchar(20), 
                            "magnitude" int, 
                            "subject" varchar(100), 
                            "group" varchar(100), 
                            "series_title_1" varchar(100), 
                            "series_title_2" varchar(100), 
                            "series_title_3" varchar(100), 
                            "series_title_4" varchar(100), 
                            "series_title_5" varchar(100), 
                            primary key: "series_reference"''',
        'myapps_organization': '''"myapps_organization" "DATABASE" 
                            "index" int, 
                            "organization_id" varchar(100), 
                            "name" varchar(255), 
                            "website" varchar(255), 
                            "country" varchar(100), 
                            "description" text, 
                            "founded" int, 
                            "industry" varchar(100), 
                            "number_of_employees" int, 
                            primary key: "index"''',
        'myapps_people': '''"myapps_people" "DATABASE" 
                            "index" int, 
                            "user_id" varchar(100), 
                            "first_name" varchar(50), 
                            "last_name" varchar(50), 
                            "sex" varchar(10), 
                            "email" varchar(100), 
                            "phone" varchar(50), 
                            "date_of_birth" date, 
                            "job_title" varchar(100), 
                            primary key: "index"''',
        'myapps_businessdata': '''"myapps_businessdata" "DATABASE" 
                            "series_reference" varchar(100), 
                            "period" varchar(10), 
                            "data_value" varchar(20), 
                            "suppressed" varchar(1), 
                            "status" varchar(1), 
                            "units" varchar(20), 
                            "magnitude" int, 
                            "subject" varchar(100), 
                            "group" varchar(100), 
                            "series_title_1" varchar(100), 
                            "series_title_2" varchar(100), 
                            "series_title_3" varchar(100), 
                            "series_title_4" varchar(100), 
                            "series_title_5" varchar(100), 
                            primary key: "series_reference"''',
        'myapps_customer': '''"myapps_customer" "DATABASE"
                        "index" int,
                        "customer_id" varchar(100),
                        "first_name" varchar(50),
                        "last_name" varchar(50),
                        "company" varchar(100),
                        "city" varchar(100),
                        "country" varchar(100),
                        "phone1" varchar(50),
                        "phone2" varchar(50),
                        "email" varchar(100),
                        "subscription_date" date,
                        "website" varchar(255),
                        primary key: "index"'''
    }

    for table in table_names:
        if table not in schema:
            conn = sqlite3.connect(db_path)
            cursor = conn.cursor()
            cursor.execute(f"PRAGMA table_info({table})")
            columns = cursor.fetchall()
            conn.close()
            column_schema = ", ".join([f'"{col[1]}" {col[2]}' for col in columns])
            schema[table] = f'"{table}" "DATABASE" {column_schema}, primary key: "{columns[0][1]}"'

    # Initialize selected_table
    selected_table = None

    # Handle CSV Upload
    upload_success = False
    if request.method == 'POST':
        upload_form = UploadCSVForm(request.POST, request.FILES)
        if upload_form.is_valid():
            csv_file = request.FILES['csv_file']
            # Handle saving the file
            file_path = os.path.join(settings.MEDIA_ROOT, csv_file.name)
            with open(file_path, 'wb+') as destination:
                for chunk in csv_file.chunks():
                    destination.write(chunk)
            
            # Assuming you want to read the CSV to get table names dynamically
            # You can modify this part based on your actual CSV structure
            try:
                df = pd.read_csv(file_path)
                table_names_from_csv = df.columns.tolist()  # Assuming first row contains headers
                table_names.extend(table_names_from_csv)  # Add to existing table_names list
            except Exception as e:
                print(f"Error reading CSV: {e}")
            
            upload_success = True
    else:
        upload_form = UploadCSVForm()

    if request.method == 'POST':
        form = QueryForm(request.POST)
        if form.is_valid():
            question = form.cleaned_data['query']
            selected_table = request.POST.get('selected_table')

            if selected_table in schema:
                sql_query = get_sql_query(question, schema[selected_table])
                results = query_database(db_path, sql_query)
                query_history.insert(0, {'question': question, 'sql_query': sql_query, 'results': results})
                # request.session['query_history'] = query_history
                if QueryHistory.objects.filter(user_id=user_id).exists():
                    QueryHistory.objects.filter(user_id=user_id).update(query_history=query_history)
                else:
                    QueryHistory.objects.create(user_id=user_id, query_history=query_history)
                if 'csv' in request.POST:
                    return download_csv(results)
                elif 'word' in request.POST:
                    return download_word(results)
                elif 'ppt' in request.POST:
                    return download_ppt(results)
                
    else:
        form = QueryForm()
        selected_table = request.GET.get('selected_table')
        if 'load_history' in request.GET:
            index = int(request.GET['load_history'])
            history_item = query_history[index]
            form = QueryForm(initial={'query': history_item['question']})
            sql_query = history_item['sql_query']
            results = history_item['results']

    context = {
        'form': form,
        'upload_form': upload_form,
        'results': results,
        'query': sql_query,
        'query_history': query_history,
        'table_names': table_names,
        'selected_table': selected_table,
        'upload_success': upload_success,
        'username': name,
        'results_json': json.dumps(results) if results else None  # Pass the results as JSON
    }
    return render(request, 'query_view.html', context)

def inventory_list(request):
    inventory_items = Inventory.objects.all()
    return render(request, 'inventory_list.html', {'inventory_items': inventory_items})
